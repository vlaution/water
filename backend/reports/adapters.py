from typing import Any
import io
import pandas as pd
from backend.reports.content import ReportContent, FormatAdapter
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet
from pptx import Presentation
from docx import Document

class PDFAdapter(FormatAdapter):
    def render(self, content: ReportContent) -> io.BytesIO:
        buffer = io.BytesIO()
        doc = SimpleDocTemplate(buffer, pagesize=letter)
        styles = getSampleStyleSheet()
        story = []
        
        story.append(Paragraph(f"Valuation Report: {content.company_name}", styles['Title']))
        story.append(Spacer(1, 12))
        
        for section in content.sections:
            story.append(Paragraph(section.title, styles['Heading1']))
            # Logic to render section.data (text, tables, charts)
            # This needs to be robust. For MVP/Architecture, we assume data has 'text'
            if "summary_text" in section.data:
                story.append(Paragraph(section.data["summary_text"], styles['Normal']))
            
            story.append(Spacer(1, 12))

        doc.build(story)
        buffer.seek(0)
        return buffer

class PPTXAdapter(FormatAdapter):
    def render(self, content: ReportContent) -> io.BytesIO:
        prs = Presentation()
        
        # Title Slide
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = f"Valuation: {content.company_name}"
        
        for section in content.sections:
             slide = prs.slides.add_slide(prs.slide_layouts[1])
             slide.shapes.title.text = section.title
             # logic to add content
        
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer

class ExcelAdapter(FormatAdapter):
    def render(self, content: ReportContent) -> io.BytesIO:
        buffer = io.BytesIO()
        with pd.ExcelWriter(buffer, engine='xlsxwriter') as writer:
            # Summary Sheet
            df_summary = pd.DataFrame([
                {"Metric": "Company", "Value": content.company_name},
                {"Metric": "Date", "Value": content.valuation_date}
            ])
            df_summary.to_excel(writer, sheet_name="Summary", index=False)
            
            # Sections
            for section in content.sections:
                # Mock logic to flatten dict to df
                df = pd.DataFrame(str(section.data), index=[0], columns=["Data"])
                df.to_excel(writer, sheet_name=section.id[:30], index=False)
                
        buffer.seek(0)
        return buffer
