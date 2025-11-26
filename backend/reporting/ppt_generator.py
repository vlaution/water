from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from io import BytesIO
from datetime import datetime

class PPTGenerator:
    def generate_presentation(self, results: dict, run_id: str) -> bytes:
        prs = Presentation()
        
        # Slide 1: Title Slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        company_name = results.get('input_summary', {}).get('company_name', 'Target Company')
        title.text = f"Valuation Report: {company_name}"
        subtitle.text = f"Generated on {datetime.now().strftime('%B %d, %Y')}\nRun ID: {run_id}"
        
        # Slide 2: Executive Summary
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = "Executive Summary"
        tf = body_shape.text_frame
        tf.text = f"Enterprise Value: ${results.get('enterprise_value', 0)/1000000:.1f}M"
        
        p = tf.add_paragraph()
        p.text = f"Equity Value: ${results.get('equity_value', 0)/1000000:.1f}M"
        
        p = tf.add_paragraph()
        p.text = f"WACC: {results.get('wacc', 0):.1%}"
        
        p = tf.add_paragraph()
        p.text = "Methodology: Discounted Cash Flow (DCF)"

        # Slide 3: Key Assumptions
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = "Key Assumptions"
        tf = body_shape.text_frame
        
        dcf_input = results.get('input_summary', {}).get('dcf_input', {}).get('projections', {})
        
        assumptions = [
            f"Revenue Growth (Start): {dcf_input.get('revenue_growth_start', 0):.1%}",
            f"EBITDA Margin (Start): {dcf_input.get('ebitda_margin_start', 0):.1%}",
            f"Terminal Growth Rate: {dcf_input.get('terminal_growth_rate', 0):.1%}",
            f"Discount Rate: {dcf_input.get('discount_rate', 0):.1%}"
        ]
        
        tf.text = assumptions[0]
        for a in assumptions[1:]:
            p = tf.add_paragraph()
            p.text = a

        # Slide 4: Sensitivity Analysis
        if 'sensitivity' in results and results['sensitivity']:
            slide = prs.slides.add_slide(prs.slide_layouts[5]) # Title Only
            title_shape = slide.shapes.title
            title_shape.text = "Sensitivity Analysis"
            
            sens = results['sensitivity']
            rows = len(sens['y_axis']['values'])
            cols = len(sens['x_axis']['values'])
            
            left = Inches(1)
            top = Inches(2)
            width = Inches(8)
            height = Inches(4)
            
            table = slide.shapes.add_table(rows + 1, cols + 1, left, top, width, height).table
            
            # Header
            table.cell(0, 0).text = f"{sens['y_axis']['name']}\\{sens['x_axis']['name']}"
            for i, val in enumerate(sens['x_axis']['values']):
                table.cell(0, i+1).text = f"{val:.1%}"
                
            # Data
            for i, row_val in enumerate(sens['y_axis']['values']):
                table.cell(i+1, 0).text = f"{row_val:.1%}"
                for j, val in enumerate(sens['matrix'][i]):
                    table.cell(i+1, j+1).text = f"${val/1000000:.1f}M"
                    table.cell(i+1, j+1).text_frame.paragraphs[0].font.size = Pt(10)

        # Save to buffer
        buffer = BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer.read()
